from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches
from reportlab.lib.utils import ImageReader
from reportlab.pdfgen import canvas


def images_to_pptx(image_paths: list[Path], output_path: Path) -> Path:
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for image_path in image_paths:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(image_path), 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(output_path)
    return output_path


def images_to_pdf(image_paths: list[Path], output_path: Path) -> Path:
    output_path.parent.mkdir(parents=True, exist_ok=True)
    if not image_paths:
        raise ValueError("no images to export")
    first = Image.open(image_paths[0])
    page_size = (first.width, first.height)
    c = canvas.Canvas(str(output_path), pagesize=page_size)
    for image_path in image_paths:
        with Image.open(image_path) as img:
            w, h = img.size
        c.setPageSize((w, h))
        c.drawImage(ImageReader(str(image_path)), 0, 0, width=w, height=h)
        c.showPage()
    c.save()
    return output_path
